#!/usr/bin/env python3
"""
File Reader Script

Usage:
    python file_reader.py --action check-env
    python file_reader.py --action locate
    python file_reader.py --action read --file-path <file_path>
    python file_reader.py --action analyze --file-path <file_path>
    python file_reader.py --action convert --file-path <file_path> --output-format <format> [--output-path <path>]
"""

import argparse
import sys
import traceback
from pathlib import Path
import os

__version__ = "1.0.0"


# ── Error Classes ──

class FileReaderError(Exception):
    """Base error with recovery suggestion."""
    def __init__(self, message: str, suggestion: str = None):
        super().__init__(message)
        self.suggestion = suggestion
    def format(self) -> str:
        msg = f"❌ Error: {self}"
        if self.suggestion:
            msg += f"\n💡 Suggestion: {self.suggestion}"
        return msg


# ── Environment Check ──

def check_environment() -> dict:
    """Check runtime dependencies. Returns status dict."""
    result = {
        "python_version": f"{sys.version_info.major}.{sys.version_info.minor}.{sys.version_info.micro}",
        "python_path": sys.executable,
        "platform": sys.platform,
        "script_path": str(Path(__file__).resolve()),
        "packages": {},
        "all_ok": True,
        "errors": [],
    }
    
    # Check Python version
    if sys.version_info < (3, 9):
        result["errors"].append(f"Python 3.9+ required")
        result["all_ok"] = False
    
    # Check packages
    required_packages = [
        "python-docx",
        "openpyxl",
        "python-pptx",
        "PyPDF2",
        "pandas"
    ]
    
    for pkg in required_packages:
        try:
            mod = __import__(pkg)
            ver = getattr(mod, "__version__", "?")
            result["packages"][pkg] = {"status": "ok", "version": str(ver)}
        except ImportError:
            result["packages"][pkg] = {"status": "missing"}
            result["errors"].append(f"Missing: {pkg}")
            result["all_ok"] = False
    
    return result


def print_env_report(info: dict):
    """Pretty-print environment check."""
    print("=" * 60)
    print("🔍 ENVIRONMENT CHECK")
    print("=" * 60)
    print(f"  Python  : {info['python_version']}")
    print(f"  Path    : {info['python_path']}")
    print(f"  Platform: {info['platform']}")
    print(f"  Script  : {info['script_path']}")
    for pkg, st in info["packages"].items():
        icon = "✅" if st["status"] == "ok" else "❌"
        ver = st.get("version", "MISSING")
        print(f"  {icon} {pkg}: {ver}")
    if info["all_ok"]:
        print("\n  ✅ All good!")
    else:
        for e in info["errors"]:
            print(f"  ❌ {e}")
    print("=" * 60)


# ── File Handlers ──

def read_text_file(file_path: Path):
    """Read text file."""
    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        lines = content.split("\n")
        print(f"📄 File: {file_path}")
        print(f"📊 Text file with {len(lines)} lines")
        print("\nContent preview:")
        print("================")
        preview = "\n".join(lines[:50])  # Show first 50 lines
        print(preview)
        if len(lines) > 50:
            print("...")
            print(f"(File truncated. Total lines: {len(lines)})")
    except Exception as e:
        raise FileReaderError(f"Failed to read text file: {e}", "Check file permissions and encoding")

def read_word_file(file_path: Path):
    """Read Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        print(f"📄 File: {file_path}")
        print(f"📊 Word document with {len(paragraphs)} paragraphs")
        print("\nContent preview:")
        print("================")
        preview = "\n\n".join(paragraphs[:5])  # Show first 5 paragraphs
        print(preview)
        if len(paragraphs) > 5:
            print("...")
            print(f"(Document truncated. Total paragraphs: {len(paragraphs)})")
    except ImportError:
        raise FileReaderError("python-docx not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to read Word file: {e}", "Check if the file is a valid Word document")

def read_excel_file(file_path: Path):
    """Read Excel file."""
    try:
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        print(f"📄 File: {file_path}")
        print(f"📊 Excel file with {len(xls.sheet_names)} sheets")
        
        for sheet_name in xls.sheet_names:
            print(f"\nSheet: {sheet_name}")
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            print(f"- Rows: {len(df)}")
            print(f"- Columns: {len(df.columns)}")
            print(f"- Columns: {', '.join(df.columns.tolist())}")
            
            if len(df) > 0:
                print("\nSample data:")
                print(df.head())
    except ImportError:
        raise FileReaderError("pandas or openpyxl not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to read Excel file: {e}", "Check if the file is a valid Excel document")

def read_powerpoint_file(file_path: Path):
    """Read PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                slides.append((i + 1, "\n".join(slide_text)))
        
        print(f"📄 File: {file_path}")
        print(f"📊 PowerPoint file with {len(slides)} slides with content")
        
        for slide_num, content in slides[:3]:  # Show first 3 slides
            print(f"\nSlide {slide_num}:")
            print("-" * 40)
            print(content)
        
        if len(slides) > 3:
            print("...")
            print(f"(Presentation truncated. Total slides with content: {len(slides)})")
    except ImportError:
        raise FileReaderError("python-pptx not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to read PowerPoint file: {e}", "Check if the file is a valid PowerPoint document")

def read_pdf_file(file_path: Path):
    """Read PDF file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = []
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            page_text = page.extract_text()
            if page_text.strip():
                text.append((page_num + 1, page_text))
        
        print(f"📄 File: {file_path}")
        print(f"📊 PDF file with {len(text)} pages with content")
        
        for page_num, content in text[:2]:  # Show first 2 pages
            print(f"\nPage {page_num}:")
            print("-" * 40)
            lines = content.split("\n")[:20]  # Show first 20 lines per page
            print("\n".join(lines))
            if len(content.split("\n")) > 20:
                print("...")
        
        if len(text) > 2:
            print("...")
            print(f"(PDF truncated. Total pages with content: {len(text)})")
    except ImportError:
        raise FileReaderError("PyPDF2 not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to read PDF file: {e}", "Check if the file is a valid PDF document")


def read_file(file_path: Path):
    """Read file based on extension."""
    ext = file_path.suffix.lower()
    
    if ext in [".txt", ".md"]:
        read_text_file(file_path)
    elif ext == ".docx":
        read_word_file(file_path)
    elif ext == ".xlsx":
        read_excel_file(file_path)
    elif ext == ".pptx":
        read_powerpoint_file(file_path)
    elif ext == ".pdf":
        read_pdf_file(file_path)
    else:
        raise FileReaderError(f"Unsupported file format: {ext}", "Check if the file type is supported")


def analyze_excel_file(file_path: Path):
    """Analyze Excel file."""
    try:
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        print(f"📄 File: {file_path}")
        print(f"📊 Excel file with {len(xls.sheet_names)} sheets")
        
        for sheet_name in xls.sheet_names:
            print(f"\nSheet: {sheet_name}")
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            print(f"- Rows: {len(df)}")
            print(f"- Columns: {len(df.columns)}")
            print(f"- Columns: {', '.join(df.columns.tolist())}")
            
            # Basic statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                print("\nBasic Statistics:")
                for col in numeric_cols:
                    if df[col].count() > 0:
                        print(f"- {col}:")
                        print(f"  Mean: {df[col].mean():.2f}")
                        print(f"  Sum: {df[col].sum():.2f}")
                        print(f"  Min: {df[col].min():.2f}")
                        print(f"  Max: {df[col].max():.2f}")
            
            # Top values for categorical columns
            categorical_cols = df.select_dtypes(include=['object']).columns
            if len(categorical_cols) > 0:
                print("\nTop values:")
                for col in categorical_cols[:2]:  # Limit to first 2 categorical columns
                    print(f"- {col}:")
                    top_values = df[col].value_counts().head(5)
                    for val, count in top_values.items():
                        print(f"  {val}: {count}")
    except ImportError:
        raise FileReaderError("pandas not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to analyze Excel file: {e}", "Check if the file is a valid Excel document")

def analyze_word_file(file_path: Path):
    """Analyze Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        words = []
        for para in paragraphs:
            words.extend(para.split())
        
        print(f"📄 File: {file_path}")
        print(f"📊 Word document analysis:")
        print(f"- Paragraphs: {len(paragraphs)}")
        print(f"- Words: {len(words)}")
        print(f"- Characters: {sum(len(p) for p in paragraphs)}")
        
        if paragraphs:
            print("\nLongest paragraph:")
            longest = max(paragraphs, key=len)
            print(longest[:200] + ("..." if len(longest) > 200 else ""))
    except ImportError:
        raise FileReaderError("python-docx not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to analyze Word file: {e}", "Check if the file is a valid Word document")

def analyze_pdf_file(file_path: Path):
    """Analyze PDF file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = []
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            page_text = page.extract_text()
            if page_text.strip():
                text.append(page_text)
        
        full_text = " ".join(text)
        words = full_text.split()
        
        print(f"📄 File: {file_path}")
        print(f"📊 PDF file analysis:")
        print(f"- Pages: {len(reader.pages)}")
        print(f"- Pages with content: {len(text)}")
        print(f"- Words: {len(words)}")
        print(f"- Characters: {len(full_text)}")
        
        if text:
            print("\nContent preview:")
            preview = text[0][:300] + ("..." if len(text[0]) > 300 else "")
            print(preview)
    except ImportError:
        raise FileReaderError("PyPDF2 not installed", "Run setup_env.py to install dependencies")
    except Exception as e:
        raise FileReaderError(f"Failed to analyze PDF file: {e}", "Check if the file is a valid PDF document")

def analyze_text_file(file_path: Path):
    """Analyze text file."""
    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        lines = content.split("\n")
        words = content.split()
        
        print(f"📄 File: {file_path}")
        print(f"📊 Text file analysis:")
        print(f"- Lines: {len(lines)}")
        print(f"- Words: {len(words)}")
        print(f"- Characters: {len(content)}")
        
        if lines:
            print("\nFirst few lines:")
            for i, line in enumerate(lines[:5]):
                print(f"{i+1}: {line}")
    except Exception as e:
        raise FileReaderError(f"Failed to analyze text file: {e}", "Check file permissions and encoding")

def analyze_file(file_path: Path):
    """Analyze file based on extension."""
    ext = file_path.suffix.lower()
    
    if ext in [".txt", ".md"]:
        analyze_text_file(file_path)
    elif ext == ".docx":
        analyze_word_file(file_path)
    elif ext == ".xlsx":
        analyze_excel_file(file_path)
    elif ext == ".pptx":
        print(f"📄 File: {file_path}")
        print("📊 PowerPoint analysis not fully supported yet")
        read_powerpoint_file(file_path)
    elif ext == ".pdf":
        analyze_pdf_file(file_path)
    else:
        raise FileReaderError(f"Unsupported file format: {ext}", "Check if the file type is supported")

def convert_to_text(file_path: Path, output_path: Path):
    """Convert file to text."""
    ext = file_path.suffix.lower()
    text_content = []
    
    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
    elif ext == ".xlsx":
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            text_content.append(f"=== Sheet: {sheet_name} ===")
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text_content.append(str(df))
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_content.append(f"=== Slide {i+1} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
    elif ext == ".pdf":
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        for page_num in range(len(reader.pages)):
            text_content.append(f"=== Page {page_num+1} ===")
            page = reader.pages[page_num]
            text_content.append(page.extract_text())
    else:
        # Already text file
        text_content.append(file_path.read_text(encoding="utf-8", errors="replace"))
    
    output_path.write_text("\n\n".join(text_content), encoding="utf-8")
    print(f"✅ Converted to text: {output_path}")

def convert_to_markdown(file_path: Path, output_path: Path):
    """Convert file to markdown."""
    ext = file_path.suffix.lower()
    md_content = []
    
    md_content.append(f"# Converted from {file_path.name}")
    md_content.append("")
    
    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                md_content.append(para.text)
                md_content.append("")
    elif ext == ".xlsx":
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            md_content.append(f"## Sheet: {sheet_name}")
            md_content.append("")
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            # Convert dataframe to markdown
            md_content.append(df.to_markdown())
            md_content.append("")
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            md_content.append(f"## Slide {i+1}")
            md_content.append("")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_content.append(shape.text)
                    md_content.append("")
    elif ext == ".pdf":
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        for page_num in range(len(reader.pages)):
            md_content.append(f"## Page {page_num+1}")
            md_content.append("")
            page = reader.pages[page_num]
            text = page.extract_text()
            if text.strip():
                md_content.append(text)
                md_content.append("")
    else:
        # Already text file, just copy
        md_content.extend(file_path.read_text(encoding="utf-8", errors="replace").split("\n"))
    
    output_path.write_text("\n".join(md_content), encoding="utf-8")
    print(f"✅ Converted to markdown: {output_path}")

def convert_file(file_path: Path, output_format: str, output_path: Path = None):
    """Convert file to specified format."""
    if not output_path:
        # Default output path
        base_name = file_path.stem
        output_path = file_path.parent / f"{base_name}.{output_format}"
    
    output_format = output_format.lower()
    
    try:
        if output_format == "txt":
            convert_to_text(file_path, output_path)
        elif output_format == "md":
            convert_to_markdown(file_path, output_path)
        elif output_format == "csv" and file_path.suffix.lower() == ".xlsx":
            import pandas as pd
            df = pd.read_excel(file_path)
            df.to_csv(output_path, index=False, encoding="utf-8-sig")
            print(f"✅ Converted to CSV: {output_path}")
        else:
            raise FileReaderError(f"Unsupported conversion: {file_path.suffix} to {output_format}", "Check supported conversion formats")
    except Exception as e:
        raise FileReaderError(f"Failed to convert file: {e}", "Check file format and permissions")


# ── CLI ──

def main():
    parser = argparse.ArgumentParser(description=f"{__doc__}")
    parser.add_argument("--action", required=True,
                        choices=["check-env", "locate", "read", "analyze", "convert"])
    parser.add_argument("--file-path", help="Path to the file to process")
    parser.add_argument("--output-format", help="Output format for conversion (txt, md, csv)")
    parser.add_argument("--output-path", help="Custom output path for converted files")
    parser.add_argument("--version", action="version",
                        version=f"%(prog)s {__version__}")

    args = parser.parse_args()

    if args.action == "locate":
        print(f"Script: {Path(__file__).resolve()}")
        print(f"Dir:    {Path(__file__).resolve().parent}")
        sys.exit(0)

    if args.action == "check-env":
        info = check_environment()
        print_env_report(info)
        sys.exit(0 if info["all_ok"] else 1)

    # Validate file path for file operations
    if args.action in ["read", "analyze", "convert"]:
        if not args.file_path:
            parser.error("--file-path is required for read, analyze, and convert actions")
        file_path = Path(args.file_path)
        if not file_path.exists():
            raise FileReaderError(f"File not found: {file_path}", "Check the file path")
        if not file_path.is_file():
            raise FileReaderError(f"Not a file: {file_path}", "Check the file path")

    # Handle file operations
    try:
        if args.action == "read":
            read_file(file_path)
        elif args.action == "analyze":
            analyze_file(file_path)
        elif args.action == "convert":
            if not args.output_format:
                parser.error("--output-format is required for convert action")
            output_path = Path(args.output_path) if args.output_path else None
            convert_file(file_path, args.output_format, output_path)
    except FileReaderError as e:
        print(e.format(), file=sys.stderr)
        sys.exit(1)
    except KeyboardInterrupt:
        print("\n⚠️ Interrupted.", file=sys.stderr)
        sys.exit(130)
    except Exception as e:
        print(f"\n❌ Unexpected error: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        print(f"\n💡 Try: python \"{__file}\" --action check-env",
              file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
